#!/usr/bin/env python3
"""
generate_pptx.py — HTR Combined Deck generator
Converts COMBINED_DECK.md to a properly designed COMBINED_DECK.pptx.

Slide design principles:
- 16:9 widescreen (13.33" × 7.5")
- Navy header bar (~1.4") with large white title (36pt)
- Large readable body text (24pt bullets, 20pt sub-bullets)
- Tables: readable row height, 13pt min cell font
- Diagrams fill available space, not a corner of it
- Speaker notes from "> " lines
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Brand colors ──────────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)

COLORS = {
    "navy":    rgb(15,  40,  80),
    "sky":     rgb(14,  165, 233),
    "emerald": rgb(16,  185, 129),
    "indigo":  rgb(99,  102, 241),
    "red":     rgb(220, 50,  50),
    "violet":  rgb(139, 92,  246),
    "teal":    rgb(20,  184, 166),
    "white":   rgb(255, 255, 255),
    "nearblack": rgb(20, 20, 40),
    "darkgray": rgb(60, 60, 80),
    "midgray": rgb(130, 130, 150),
    "lightgray": rgb(220, 222, 228),
    "offwhite": rgb(248, 249, 252),
    "rowalt":  rgb(240, 244, 252),
}

PILLAR_COLORS = {
    "policy":     "sky",
    "economics":  "emerald",
    "technology": "indigo",
    "clinical":   "red",
    "equity":     "violet",
    "operations": "teal",
}

STEP_COLORS = ["navy", "sky", "emerald", "indigo", "red", "violet", "teal"]

# ── Slide dimensions ──────────────────────────────────────────────────────────
W          = Inches(13.333)
H          = Inches(7.5)
HDR_H      = Inches(1.45)      # header bar height
MARGIN_L   = Inches(0.55)
MARGIN_R   = Inches(0.55)
BODY_TOP   = HDR_H + Inches(0.22)
BODY_W     = W - MARGIN_L - MARGIN_R
BODY_BOT   = H - Inches(0.18)
BODY_H     = BODY_BOT - BODY_TOP
ACCENT_W   = Inches(0.1)       # left accent strip width


def c(name):
    return COLORS.get(name, COLORS["navy"])


# ── Parse annotation block ────────────────────────────────────────────────────
def parse_annotation(line):
    """Parse [KIND:key=val|key=val] into (KIND, {key: val})."""
    m = re.match(r'^\[(\w+):(.+)\]$', line.strip())
    if not m:
        return None
    kind = m.group(1).upper()
    raw = m.group(2)
    fields = {}
    for part in raw.split('|'):
        if '=' in part:
            k, v = part.split('=', 1)
            fields[k.strip()] = v.strip()
    return kind, fields


# ── Markdown helpers ──────────────────────────────────────────────────────────
def strip_md(text):
    """Strip markdown bold/italic markers, return plain text."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    text = re.sub(r'_(.+?)_', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    return text


def is_bold_md(text):
    return bool(re.search(r'\*\*', text) or re.search(r'__', text))


# ── Header ────────────────────────────────────────────────────────────────────
def add_header(slide, title_text, accent_color="sky"):
    # Full-width navy bar
    bar = slide.shapes.add_shape(1, 0, 0, W, HDR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = c("navy")
    bar.line.fill.background()

    # Left accent strip (pillar color)
    strip = slide.shapes.add_shape(1, 0, 0, ACCENT_W, HDR_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = c(accent_color)
    strip.line.fill.background()

    # Title text box (inset from strip)
    txt = slide.shapes.add_textbox(ACCENT_W + Inches(0.3), Inches(0.2),
                                   W - ACCENT_W - Inches(0.4), HDR_H - Inches(0.3))
    tf = txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = strip_md(title_text)
    run.font.bold = True
    run.font.size = Pt(34)
    run.font.color.rgb = c("white")
    run.font.name = "Calibri"


# ── Body text ─────────────────────────────────────────────────────────────────
def add_body_text(slide, lines, top=None, left=None, width=None, height=None,
                  base_size=24):
    """Render a list of text lines as a formatted text box."""
    top    = top    if top    is not None else BODY_TOP
    left   = left   if left   is not None else MARGIN_L
    width  = width  if width  is not None else BODY_W
    height = height if height is not None else BODY_H

    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for line in lines:
        raw = line.rstrip()
        if not raw.strip():
            # blank line → small spacer paragraph
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(4)
            continue

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        stripped = raw.lstrip()

        # Numbered list item
        num_m = re.match(r'^(\d+)\.\s+(.*)', stripped)
        # Bullet item
        is_bullet = stripped.startswith(('- ', '* ', '• '))
        # Sub-bullet (indented)
        indent_level = len(raw) - len(raw.lstrip())
        is_sub = indent_level >= 2 and is_bullet

        if num_m:
            label = num_m.group(1) + '.  '
            content = num_m.group(2)
            font_size = base_size
            p.space_before = Pt(6)
        elif is_bullet:
            label = '•  '
            content = stripped[2:]
            font_size = base_size if not is_sub else max(base_size - 4, 18)
            p.space_before = Pt(4)
        else:
            label = None
            content = stripped
            # Heading-style lines (### or bold-only)
            if stripped.startswith('### '):
                content = stripped[4:]
                font_size = base_size + 2
            elif stripped.startswith('## '):
                content = stripped[3:]
                font_size = base_size + 4
            elif stripped.startswith('# '):
                content = stripped[2:]
                font_size = base_size + 6
            else:
                font_size = base_size
            p.space_before = Pt(8)

        bold = is_bold_md(content) or stripped.startswith(('## ', '### ', '**'))
        content_clean = strip_md(content)

        p.alignment = PP_ALIGN.LEFT

        if label:
            r0 = p.add_run()
            r0.text = label
            r0.font.size = Pt(font_size)
            r0.font.color.rgb = c("navy")
            r0.font.bold = True
            r0.font.name = "Calibri"

        run = p.add_run()
        run.text = content_clean
        run.font.size = Pt(font_size)
        run.font.color.rgb = c("nearblack")
        run.font.bold = bold
        run.font.name = "Calibri"


# ── TABLE renderer ────────────────────────────────────────────────────────────
def render_table(slide, fields, top, left, avail_h, avail_w=None):
    avail_w = avail_w or BODY_W
    title   = fields.get('title', '')
    headers = [h.strip() for h in fields.get('headers', '').split(',') if h.strip()]
    raw_rows = fields.get('rows', '')

    rows_data = []
    for rs in raw_rows.split('|'):
        rs = rs.strip()
        if rs:
            rows_data.append([c2.strip() for c2 in rs.split(';')])

    if not headers and not rows_data:
        return top

    # Optional sub-title
    if title:
        tb = slide.shapes.add_textbox(left, top, avail_w, Inches(0.38))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = c("navy")
        run.font.name = "Calibri"
        top += Inches(0.42)
        avail_h -= Inches(0.42)

    ncols = max(len(headers), max((len(r) for r in rows_data), default=1))
    nrows_data = len(rows_data)
    nrows_total = nrows_data + (1 if headers else 0)

    if nrows_total == 0 or ncols == 0:
        return top

    # Row heights: header 0.55", data rows fill remaining space, min 0.48"
    hdr_h  = Inches(0.55)
    data_h = max(Inches(0.48), (avail_h - hdr_h) / nrows_data) if nrows_data else Inches(0.5)
    # Cap data row height so table doesn't overflow
    data_h = min(data_h, Inches(0.72))
    tbl_h = hdr_h + data_h * nrows_data

    tbl_w = avail_w
    tbl = slide.shapes.add_table(nrows_total, ncols, left, top, int(tbl_w), int(tbl_h)).table

    col_w = tbl_w // ncols
    for ci in range(ncols):
        tbl.columns[ci].width = int(col_w)

    # Header row
    if headers:
        for ci, htext in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = c("navy")
            p2 = cell.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run = p2.add_run()
            run.text = htext
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = c("white")
            run.font.name = "Calibri"

    # Data rows
    for ri, row in enumerate(rows_data):
        tr = ri + (1 if headers else 0)
        for ci in range(ncols):
            cell = tbl.cell(tr, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = c("rowalt") if ri % 2 == 0 else c("white")
            val = row[ci] if ci < len(row) else ''
            p2 = cell.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            run = p2.add_run()
            run.text = strip_md(val)
            run.font.size = Pt(13)
            run.font.bold = ci == 0  # bold first column
            run.font.color.rgb = c("nearblack")
            run.font.name = "Calibri"

    return top + tbl_h + Inches(0.15)


# ── DATAFLOW renderer ─────────────────────────────────────────────────────────
def render_dataflow(slide, fields, top, left, avail_h, avail_w=None):
    avail_w = avail_w or BODY_W
    title = fields.get('title', '')
    steps = [s.strip() for s in fields.get('steps', '').split('|') if s.strip()]
    if not steps:
        return top

    if title:
        tb = slide.shapes.add_textbox(left, top, avail_w, Inches(0.38))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = c("navy")
        run.font.name = "Calibri"
        top += Inches(0.44)
        avail_h -= Inches(0.44)

    n = len(steps)
    arrow_w = Inches(0.3)
    box_h = min(Inches(1.4), avail_h - Inches(0.1))
    total_arrow = arrow_w * (n - 1)
    box_w = (avail_w - total_arrow) / n

    for i, step in enumerate(steps):
        bx = left + i * (box_w + arrow_w)
        by = top
        color_name = STEP_COLORS[i % len(STEP_COLORS)]

        box = slide.shapes.add_shape(5, int(bx), int(by), int(box_w), int(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = c(color_name)
        box.line.color.rgb = rgb(200, 200, 210)
        box.line.width = Pt(0.5)

        tf2 = box.text_frame
        tf2.word_wrap = True
        tf2.margin_left  = Inches(0.12)
        tf2.margin_right = Inches(0.12)
        tf2.margin_top   = Inches(0.08)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER

        # Step number label
        num_run = p2.add_run()
        num_run.text = f"{i+1}\n"
        num_run.font.size = Pt(18)
        num_run.font.bold = True
        num_run.font.color.rgb = rgb(255, 255, 255)
        num_run.font.name = "Calibri"

        p3 = tf2.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        txt_run = p3.add_run()
        txt_run.text = step
        txt_run.font.size = Pt(13)
        txt_run.font.bold = False
        txt_run.font.color.rgb = rgb(230, 240, 255)
        txt_run.font.name = "Calibri"

        # Arrow (except after last)
        if i < n - 1:
            ax = bx + box_w
            ay = by + box_h / 2 - Inches(0.18)
            arr = slide.shapes.add_shape(1, int(ax), int(ay), int(arrow_w), int(Inches(0.36)))
            arr.fill.solid()
            arr.fill.fore_color.rgb = c("lightgray")
            arr.line.fill.background()
            # Arrow text
            atf = arr.text_frame
            arun = atf.paragraphs[0].add_run()
            arun.text = "▶"
            arun.font.size = Pt(18)
            arun.font.color.rgb = c("midgray")
            atf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return top + box_h + Inches(0.15)


# ── DIAGRAM / FLOWCHART renderer ──────────────────────────────────────────────
def render_diagram_flowchart(slide, fields, top, left, avail_h, avail_w=None):
    avail_w = avail_w or BODY_W
    title = fields.get('title', '')
    raw   = fields.get('data', '')
    color_names = [cn.strip() for cn in fields.get('colors', '').split(',') if cn.strip()]

    # Parse nodes (unique, in order) and edges
    edges = []
    nodes_ordered = []
    seen = set()
    for token in re.split(r'[|]', raw):
        token = token.strip()
        if '->' in token:
            parts = [p.strip() for p in token.split('->')]
            for p in parts:
                if p and p not in seen:
                    nodes_ordered.append(p)
                    seen.add(p)
            for j in range(len(parts) - 1):
                edges.append((parts[j], parts[j+1]))
        elif token and token not in seen:
            nodes_ordered.append(token)
            seen.add(token)

    if not nodes_ordered:
        return top

    if title:
        tb = slide.shapes.add_textbox(left, top, avail_w, Inches(0.38))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = c("navy")
        run.font.name = "Calibri"
        top += Inches(0.44)
        avail_h -= Inches(0.44)

    n = len(nodes_ordered)
    gap = Inches(0.28)
    box_h = min(Inches(1.2), avail_h - Inches(0.1))
    box_w = min(Inches(2.6), (avail_w - gap * (n - 1)) / n)
    node_pos = {}

    for i, node in enumerate(nodes_ordered):
        col_name = color_names[i] if i < len(color_names) else STEP_COLORS[i % len(STEP_COLORS)]
        bx = left + i * (box_w + gap)
        by = top
        node_pos[node] = (bx + box_w / 2, by + box_h / 2)

        box = slide.shapes.add_shape(5, int(bx), int(by), int(box_w), int(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = c(col_name)
        box.line.color.rgb = rgb(180, 190, 210)
        box.line.width = Pt(0.5)

        tf2 = box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_right = Inches(0.1)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = node
        run2.font.size = Pt(13)
        run2.font.bold = True
        run2.font.color.rgb = c("white")
        run2.font.name = "Calibri"

    # Arrows between connected nodes
    for (src, dst) in edges:
        if src in node_pos and dst in node_pos:
            sx, sy = node_pos[src]
            dx, dy = node_pos[dst]
            if abs(dx - sx) > 50:
                ax = min(sx, dx)
                ay = sy - Inches(0.06)
                aw = abs(dx - sx)
                ah = Inches(0.12)
                arr = slide.shapes.add_shape(1, int(ax), int(ay), int(aw), int(ah))
                arr.fill.solid()
                arr.fill.fore_color.rgb = c("lightgray")
                arr.line.fill.background()

    return top + box_h + Inches(0.15)


# ── MATRIX renderer ───────────────────────────────────────────────────────────
def render_matrix(slide, fields, top, left, avail_h, avail_w=None):
    avail_w = avail_w or BODY_W
    title    = fields.get('title', '')
    row_lbls = [r.strip() for r in fields.get('rows', '').split(',') if r.strip()]
    col_lbls = [c2.strip() for c2 in fields.get('cols', '').split(',') if c2.strip()]
    raw_cells = fields.get('cells', '')

    cell_rows = []
    for row_str in raw_cells.split(';'):
        row_str = row_str.strip()
        if row_str:
            cell_rows.append([v.strip() for v in row_str.split(',')])

    nrows = len(row_lbls) + 1
    ncols = len(col_lbls) + 1
    if nrows <= 1 or ncols <= 1:
        return top

    if title:
        tb = slide.shapes.add_textbox(left, top, avail_w, Inches(0.38))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = c("navy")
        run.font.name = "Calibri"
        top += Inches(0.44)
        avail_h -= Inches(0.44)

    row_h = max(Inches(0.44), min(Inches(0.65), avail_h / nrows))
    tbl_h = row_h * nrows
    tbl_w = avail_w

    tbl = slide.shapes.add_table(nrows, ncols, left, top, int(tbl_w), int(tbl_h)).table

    label_w = Inches(1.6)
    data_col_w = (tbl_w - label_w) // (ncols - 1) if ncols > 1 else tbl_w
    tbl.columns[0].width = int(label_w)
    for ci in range(1, ncols):
        tbl.columns[ci].width = int(data_col_w)

    # Corner
    corner = tbl.cell(0, 0)
    corner.fill.solid()
    corner.fill.fore_color.rgb = c("navy")

    # Column headers
    for ci, clbl in enumerate(col_lbls):
        cell = tbl.cell(0, ci + 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = c("navy")
        p2 = cell.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run = p2.add_run()
        run.text = clbl
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = c("white")
        run.font.name = "Calibri"

    # Row labels + data cells
    for ri, rlbl in enumerate(row_lbls):
        lc = tbl.cell(ri + 1, 0)
        lc.fill.solid()
        lc.fill.fore_color.rgb = c("navy")
        run = lc.text_frame.paragraphs[0].add_run()
        run.text = rlbl
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = c("white")
        run.font.name = "Calibri"

        for ci in range(len(col_lbls)):
            dc = tbl.cell(ri + 1, ci + 1)
            val = ''
            if ri < len(cell_rows) and ci < len(cell_rows[ri]):
                val = cell_rows[ri][ci]

            if val == '—' or ri == ci:
                dc.fill.solid()
                dc.fill.fore_color.rgb = rgb(200, 205, 215)
                font_col = c("midgray")
            elif val.lower() in ('yes', 'true', 'x', '1'):
                dc.fill.solid()
                dc.fill.fore_color.rgb = c("emerald")
                font_col = c("white")
            else:
                dc.fill.solid()
                dc.fill.fore_color.rgb = c("offwhite")
                font_col = c("darkgray")

            p2 = dc.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = val
            run2.font.size = Pt(13)
            run2.font.bold = val.lower() in ('yes', 'x', '1')
            run2.font.color.rgb = font_col
            run2.font.name = "Calibri"

    return top + tbl_h + Inches(0.15)


# ── BAR CHART renderer ────────────────────────────────────────────────────────
def render_chart_bar(slide, fields, top, left, avail_h, avail_w=None):
    avail_w = avail_w or BODY_W
    title  = fields.get('title', 'Chart')
    labels = [l.strip() for l in fields.get('labels', '').split(',') if l.strip()]
    values = []
    for v in fields.get('values', '').split(','):
        try:
            values.append(float(v.strip()))
        except ValueError:
            values.append(0.0)
    color_names = [cn.strip() for cn in fields.get('colors', 'navy').split(',')]

    if not labels or not values:
        return top

    chart_data = ChartData()
    chart_data.categories = labels
    chart_data.add_series('', values)

    chart_h = min(avail_h - Inches(0.1), Inches(4.8))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, int(avail_w), int(chart_h), chart_data
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    title_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(16)
    title_run.font.color.rgb = c("navy")
    title_run.font.bold = True
    chart.has_legend = False

    series = chart.series[0]
    for i, pt in enumerate(series.points):
        col_name = color_names[i % len(color_names)]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c(col_name)

    # Category axis font
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.size = Pt(12)

    return top + chart_h + Inches(0.15)


# ── Slide builder ─────────────────────────────────────────────────────────────
def build_slide(prs, raw_text):
    lines = raw_text.strip().split('\n')
    if not lines:
        return

    # Extract title from first # heading
    title_text = "Slide"
    accent_color = "sky"
    body_lines = []
    notes_lines = []
    annotation_blocks = []
    found_title = False

    for line in lines:
        stripped = line.strip()

        if stripped.startswith('> '):
            notes_lines.append(stripped[2:])
            continue

        ann = parse_annotation(stripped)
        if ann:
            annotation_blocks.append(ann)
            continue

        if not found_title and stripped.startswith('# '):
            title_text = stripped[2:].strip()
            found_title = True
            # Pick accent color from title keywords
            tl = title_text.lower()
            for pname, pcolor in PILLAR_COLORS.items():
                if pname in tl:
                    accent_color = pcolor
                    break
            continue

        body_lines.append(line)

    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = c("white")

    add_header(slide, title_text, accent_color)

    # Clean body lines (remove leading/trailing blanks)
    while body_lines and not body_lines[0].strip():
        body_lines.pop(0)
    while body_lines and not body_lines[-1].strip():
        body_lines.pop()

    has_annotations = len(annotation_blocks) > 0
    has_body = bool(body_lines)

    cur_top = BODY_TOP
    avail_h = BODY_H

    if has_body and has_annotations:
        # Body text in upper portion, annotation fills lower portion
        # Estimate body height by line count
        n_lines = len([l for l in body_lines if l.strip()])
        body_h = min(Inches(0.38) * n_lines + Inches(0.2), avail_h * 0.42)
        add_body_text(slide, body_lines, top=cur_top, height=body_h, base_size=20)
        cur_top += body_h + Inches(0.12)
        avail_h -= body_h + Inches(0.12)

    elif has_body and not has_annotations:
        # Body text fills slide — pick font size based on content length
        total_chars = sum(len(l) for l in body_lines)
        n_lines = len([l for l in body_lines if l.strip()])
        if n_lines <= 4 and total_chars < 300:
            base_size = 28
        elif n_lines <= 8:
            base_size = 24
        else:
            base_size = 20
        add_body_text(slide, body_lines, top=cur_top, height=avail_h, base_size=base_size)

    # Render annotation blocks — divide available space equally
    n_anno = len(annotation_blocks)
    if n_anno > 0:
        per_h = avail_h / n_anno
        for idx, (kind, fields) in enumerate(annotation_blocks):
            item_top = cur_top + idx * per_h
            item_avail = per_h - Inches(0.08)
            kw = dict(top=int(item_top), left=int(MARGIN_L),
                      avail_h=int(item_avail), avail_w=int(BODY_W))
            if kind == 'TABLE':
                render_table(slide, fields, **kw)
            elif kind == 'CHART':
                render_chart_bar(slide, fields, **kw)
            elif kind == 'DATAFLOW':
                render_dataflow(slide, fields, **kw)
            elif kind == 'DIAGRAM':
                render_diagram_flowchart(slide, fields, **kw)
            elif kind == 'MATRIX':
                render_matrix(slide, fields, **kw)

    # Speaker notes
    if notes_lines:
        slide.notes_slide.notes_text_frame.text = ' '.join(notes_lines)

    return slide


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    src  = '/Users/baba/Vermont-Health-Platform/COMBINED_DECK.md'
    dest = '/Users/baba/Vermont-Health-Platform/COMBINED_DECK.pptx'

    with open(src, encoding='utf-8') as f:
        content = f.read()

    raw_slides = re.split(r'\n---\n', content)

    slide_blocks = []
    for block in raw_slides:
        b = block.strip()
        if not b:
            continue
        if re.search(r'^# Slide', b, re.MULTILINE) or re.search(r'\[(TABLE|CHART|DATAFLOW|DIAGRAM|MATRIX):', b):
            slide_blocks.append(b)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for block in slide_blocks:
        build_slide(prs, block)

    prs.save(dest)
    size_kb = os.path.getsize(dest) // 1024
    n_slides = len(list(prs.slides))
    print(f"Generated {n_slides} slides → {dest} ({size_kb} KB)")


if __name__ == '__main__':
    main()
