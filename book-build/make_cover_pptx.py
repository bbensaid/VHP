#!/usr/bin/env python3
"""
make_cover_pptx.py — the six-pillar framework as an EDITABLE PowerPoint.

Every circle, arrow and label is a real PowerPoint shape. Open it, drag
anything, retype anything, restyle anything. Nothing is baked into pixels.

WHY THIS EXISTS
---------------
The supplied cover art had correct label TEXT but attached several labels to the
wrong arrows — "Enables mandatory authority" drawn between Economics and
Technology when it is Policy -> Economics, "Enables risk stratification" on
Technology -> Economics when it is Technology -> Clinical, and others. A raster
image cannot be corrected; the geometry is fixed.

Here the geometry is generated FROM the book. Each of the fifteen dependencies
in Figure 1.3 is drawn as a connector from its true source pillar to its true
target pillar, and its label is placed at that connector's midpoint. The arrow
and the words cannot disagree, because both come from the same row of the table.

    python3 book-build/make_cover_pptx.py
"""
import math
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
BOOK = os.path.join(HERE, '..', 'HTR_Book_v42.md')
OUT = os.path.join(HERE, '..', 'Six-Pillar_Framework_EDITABLE.pptx')

# ── read the six-pillar table and Figure 1.3 straight from the manuscript ────
md = open(BOOK, encoding='utf-8').read()
lines = md.split('\n')

QUESTION = {}
for m in re.finditer(r'^\|\s*\*\*(\w+)\*\*\s*\|\s*(Is it [^|]+?)\s*\|', md, re.M):
    name = m.group(1).upper()
    if name not in QUESTION:
        QUESTION[name] = m.group(2).strip()

i = next(k for k, l in enumerate(lines) if l.startswith('| From ↓ / To → |'))
cols = [c.strip().upper() for c in lines[i].strip('|').split('|')][1:]
DEPS = []                      # (source, target, verb, payload)
for r in range(i + 2, i + 8):
    cells = [c.strip() for c in lines[r].strip('|').split('|')]
    src = re.sub(r'\*', '', cells[0]).strip().upper()
    for col, val in zip(cols, cells[1:]):
        if val == '—':
            continue
        m = re.match(r'\[(\w+)\]\s*(.+)', val)
        if m:
            DEPS.append((src, col, m.group(1), m.group(2).strip()))

if len(DEPS) != 15:
    sys.exit(f"expected 15 dependencies from Figure 1.3, parsed {len(DEPS)}")

PILLARS = ['POLICY', 'TECHNOLOGY', 'ECONOMICS', 'CLINICAL', 'EQUITY', 'OPERATIONS']
PCOLOR = {
    'POLICY':     RGBColor(0x1F, 0x6F, 0xB8),
    'TECHNOLOGY': RGBColor(0xE8, 0x8B, 0x1B),
    'ECONOMICS':  RGBColor(0x1E, 0x8E, 0x4E),
    'CLINICAL':   RGBColor(0x7A, 0x3E, 0xA8),
    'EQUITY':     RGBColor(0xC2, 0x2E, 0x8F),
    'OPERATIONS': RGBColor(0x8E, 0x1B, 0x2A),
}
NAVY = RGBColor(0x1B, 0x3A, 0x6B)
INK = RGBColor(0x1A, 0x1A, 0x1A)

# ── slide: 13.333 x 7.5in widescreen ────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])       # blank

SW, SH = 13.333, 7.5
CX, CY = SW / 2, SH / 2 + 0.15
RING_R = 2.05
NODE_D = 1.55                                            # circle diameter

pos = {}
for n, name in enumerate(PILLARS):
    a = -math.pi / 2 + n * (2 * math.pi / 6)
    pos[name] = (CX + RING_R * math.cos(a), CY + RING_R * math.sin(a))


def add_circle(name):
    x, y = pos[name]
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - NODE_D / 2), Inches(y - NODE_D / 2),
        Inches(NODE_D), Inches(NODE_D))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shp.line.color.rgb = PCOLOR[name]
    shp.line.width = Pt(3)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name + ':'
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = PCOLOR[name]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = QUESTION.get(name, '')
    r2.font.size = Pt(11)
    r2.font.color.rgb = INK
    return shp


circles = {n: add_circle(n) for n in PILLARS}


def edge_points(a, b, spread=0.0):
    """Rim-to-rim points between two circles, offset sideways so that two
    dependencies running between the same pair do not overlap."""
    ax, ay = pos[a]
    bx, by = pos[b]
    ang = math.atan2(by - ay, bx - ax)
    px, py = -math.sin(ang), math.cos(ang)               # perpendicular
    r = NODE_D / 2 + 0.04
    return ((ax + r * math.cos(ang) + spread * px,
             ay + r * math.sin(ang) + spread * py),
            (bx - r * math.cos(ang) + spread * px,
             by - r * math.sin(ang) + spread * py))


# how many dependencies share each unordered pair, so they can be fanned out
from collections import defaultdict
pair_n = defaultdict(int)
for s, t, _, _ in DEPS:
    pair_n[frozenset((s, t))] += 1
pair_i = defaultdict(int)

for src, tgt, verb, payload in DEPS:
    key = frozenset((src, tgt))
    k = pair_i[key]
    pair_i[key] += 1
    spread = 0.0 if pair_n[key] == 1 else (0.16 if k == 0 else -0.16)

    (x1, y1), (x2, y2) = edge_points(src, tgt, spread)
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = PCOLOR[src]
    conn.line.width = Pt(2.25)
    # arrowhead at the target end
    ln = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn as _qn
    from lxml import etree
    tail = etree.SubElement(ln, _qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')

    # Put the label ON the connector, not in a floating textbox.
    #
    # A separate textbox has to be positioned near "its" arrow, and proximity is
    # not ownership: arrows that cross the ring's interior pass close to other
    # arrows' labels, which is exactly how the supplied art ended up with words
    # attached to the wrong relationship. A connector in PowerPoint can carry
    # its own text, so here the words are part of the arrow object. Drag the
    # arrow and the label goes with it; they cannot be separated or mismatched.
    # python-pptx does not expose a text frame on a connector, but the format
    # allows one: add <p:txBody> to the <p:cxnSp> directly.
    from pptx.oxml.ns import qn as _q
    from pptx.text.text import TextFrame
    _tx = etree.SubElement(conn._element, _q('p:txBody'))
    etree.SubElement(_tx, _q('a:bodyPr')).set('wrap', 'square')
    etree.SubElement(_tx, _q('a:lstStyle'))
    etree.SubElement(_tx, _q('a:p'))
    tf = TextFrame(_tx, conn)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{verb.title()} {payload}"
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = PCOLOR[src]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"{src.title()} \u2192 {tgt.title()}"
    r2.font.size = Pt(7)
    r2.font.color.rgb = RGBColor(0x70, 0x78, 0x84)

prs.save(OUT)
print(f"wrote {OUT}")
print(f"  {len(PILLARS)} pillars, {len(DEPS)} dependencies — every shape editable")
print("  each label sits on the connector for its own row of Figure 1.3")
