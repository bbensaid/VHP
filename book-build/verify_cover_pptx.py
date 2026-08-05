#!/usr/bin/env python3
"""
verify_cover_pptx.py — check each arrow's GEOMETRY, not just its words.

The supplied cover art failed a way a text-only check could never catch: the
label wording was right, but several labels sat on the WRONG arrow — "Enables
mandatory authority" drawn between Economics and Technology when the book says
Policy → Economics, and others. Comparing word lists passed it. Comparing
geometry does not.

This reads the generated .pptx back and, for every connector, works out which
pillar circle it actually starts at and which it actually ends at, then checks
that pair — and the label the connector carries — against Figure 1.3.

Because each label lives INSIDE its connector (<p:txBody> within <p:cxnSp>),
the arrow and its words are one object and cannot be mismatched. This script
proves that holds for all fifteen.

Exit 0 = every arrow points where the book says it should.

    python3 book-build/verify_cover_pptx.py
"""
import math
import os
import re
import sys

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
BOOK = os.path.join(HERE, '..', 'HTR_Book_v42.md')
PPTX = os.path.join(HERE, '..', 'Six-Pillar_Framework_EDITABLE.pptx')
A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
E = 914400.0
NAMES = {'POLICY', 'TECHNOLOGY', 'ECONOMICS', 'CLINICAL', 'EQUITY', 'OPERATIONS'}

# ── Figure 1.3, straight from the manuscript ────────────────────────────────
lines = open(BOOK, encoding='utf-8').read().split('\n')
i = next(k for k, l in enumerate(lines) if l.startswith('| From ↓ / To → |'))
cols = [c.strip().upper() for c in lines[i].strip('|').split('|')][1:]
book = {}
for r in range(i + 2, i + 8):
    cells = [c.strip() for c in lines[r].strip('|').split('|')]
    src = re.sub(r'\*', '', cells[0]).strip().upper()
    for col, val in zip(cols, cells[1:]):
        if val == '—':
            continue
        m = re.match(r'\[(\w+)\]\s*(.+)', val)
        if m:
            book[(src, col)] = (m.group(1).upper(), m.group(2).strip())

prs = Presentation(PPTX)
slide = prs.slides[0]

# ── the six circles ─────────────────────────────────────────────────────────
centres = {}
for sh in slide.shapes:
    if not sh.has_text_frame or sh.width != sh.height:
        continue
    head = sh.text_frame.text.strip().split(':')[0].strip().upper()
    if head in NAMES:
        centres[head] = ((sh.left + sh.width / 2) / E,
                         (sh.top + sh.height / 2) / E)


def nearest(pt):
    return min(centres, key=lambda n: math.hypot(centres[n][0] - pt[0],
                                                 centres[n][1] - pt[1]))


# ── every connector, with the label it carries ──────────────────────────────
arrows = []
for sh in slide.shapes:
    if not sh.element.tag.endswith('}cxnSp'):
        continue
    x1, y1 = sh.left / E, sh.top / E
    x2, y2 = (sh.left + sh.width) / E, (sh.top + sh.height) / E
    xf = sh.element.find(f'.//{A}xfrm')
    fh = xf is not None and xf.get('flipH') == '1'
    fv = xf is not None and xf.get('flipV') == '1'
    start = ((x2 if fh else x1), (y2 if fv else y1))
    end = ((x1 if fh else x2), (y1 if fv else y2))
    text = ' '.join(t.strip() for t in sh.element.itertext() if t.strip())
    arrows.append((start, end, text))

fail = []
print(f"circles: {len(centres)}   connectors: {len(arrows)}\n")
if len(centres) != 6:
    fail.append(f"expected 6 pillar circles, found {len(centres)}")
if len(arrows) != 15:
    fail.append(f"expected 15 connectors, found {len(arrows)}")

drawn = {}
for start, end, text in arrows:
    a, b = nearest(start), nearest(end)
    drawn[(a, b)] = text

    real = book.get((a, b))
    if real is None:
        fail.append(f"an arrow runs {a}->{b}, which is not a dependency in Figure 1.3")
        continue
    verb, payload = real
    # the connector's own text must name this relationship, in the book's words
    want = f"{verb.title()} {payload}".lower()
    if want not in text.lower():
        fail.append(f"{a}->{b} arrow carries {text[:60]!r}, book says {want!r}")
    # and its direction caption must agree with where it is actually drawn
    m = re.search(r'(\w+)\s*→\s*(\w+)', text)
    if not m:
        fail.append(f"{a}->{b} arrow has no direction caption")
    elif (m.group(1).upper(), m.group(2).upper()) != (a, b):
        fail.append(f"an arrow drawn {a}->{b} is captioned "
                    f"{m.group(1)}->{m.group(2)}")

for key in book:
    if key not in drawn:
        fail.append(f"Figure 1.3 has {key[0]}->{key[1]} but nothing is drawn that way")

if fail:
    print(f"FAILED — {len(fail)} problem(s):\n")
    for f in fail:
        print("  •", f)
    sys.exit(1)

print("All 15 arrows verified against Figure 1.3:")
print("  • each runs between the correct two pillars, in the correct direction")
print("  • each carries the book's own verb and payload")
print("  • label and arrow are one object, so they cannot be mismatched")
